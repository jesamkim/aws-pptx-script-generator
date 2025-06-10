"""Validation Utilities.

This module provides comprehensive validation functions for user inputs,
file formats, and data quality assurance.
"""

import re
from typing import Dict, Any, List, Optional, Union
from pathlib import Path
from loguru import logger


class ValidationResult:
    """Represents the result of a validation operation."""
    
    def __init__(self, is_valid: bool, message: str = "", errors: List[str] = None):
        """Initialize validation result.
        
        Args:
            is_valid: Whether validation passed
            message: Validation message
            errors: List of validation errors
        """
        self.is_valid = is_valid
        self.message = message
        self.errors = errors or []
    
    def __bool__(self):
        """Return validation status."""
        return self.is_valid


class PersonaValidator:
    """Validates SA persona input data."""
    
    @staticmethod
    def validate_persona_data(persona_data: Dict[str, Any]) -> ValidationResult:
        """Validate SA persona data.
        
        Args:
            persona_data: Dictionary containing persona information
            
        Returns:
            ValidationResult object
        """
        errors = []
        
        # Required fields
        required_fields = ['full_name', 'job_title']
        for field in required_fields:
            if not persona_data.get(field, '').strip():
                errors.append(f"{field.replace('_', ' ').title()} is required")
        
        # Name validation
        full_name = persona_data.get('full_name', '').strip()
        if full_name and len(full_name) < 2:
            errors.append("Full name must be at least 2 characters")
        if full_name and len(full_name) > 100:
            errors.append("Full name must be less than 100 characters")
        
        # Job title validation
        job_title = persona_data.get('job_title', '').strip()
        if job_title and len(job_title) > 200:
            errors.append("Job title must be less than 200 characters")
        
        # Valid presentation styles
        valid_styles = ["Technical Deep-dive", "Executive Overview", "Customer Demo", "Training Session"]
        presentation_style = persona_data.get('presentation_style')
        if presentation_style and presentation_style not in valid_styles:
            errors.append(f"Invalid presentation style. Must be one of: {', '.join(valid_styles)}")
        
        # Valid experience levels
        valid_levels = ["Junior", "Senior", "Principal", "Distinguished"]
        experience_level = persona_data.get('experience_level')
        if experience_level and experience_level not in valid_levels:
            errors.append(f"Invalid experience level. Must be one of: {', '.join(valid_levels)}")
        
        is_valid = len(errors) == 0
        message = "Persona data is valid" if is_valid else f"Found {len(errors)} validation errors"
        
        return ValidationResult(is_valid, message, errors)


class PresentationValidator:
    """Validates presentation parameter data."""
    
    @staticmethod
    def validate_presentation_params(params: Dict[str, Any]) -> ValidationResult:
        """Validate presentation parameters.
        
        Args:
            params: Dictionary containing presentation parameters
            
        Returns:
            ValidationResult object
        """
        errors = []
        
        # Duration validation
        duration = params.get('duration')
        if duration is not None:
            if not isinstance(duration, (int, float)):
                errors.append("Duration must be a number")
            elif duration < 5:
                errors.append("Duration must be at least 5 minutes")
            elif duration > 120:
                errors.append("Duration must be no more than 120 minutes")
        
        # Target audience validation
        valid_audiences = ["Technical", "Executive", "Mixed", "Customer", "Internal"]
        target_audience = params.get('target_audience')
        if target_audience and target_audience not in valid_audiences:
            errors.append(f"Invalid target audience. Must be one of: {', '.join(valid_audiences)}")
        
        # Technical depth validation
        technical_depth = params.get('technical_depth')
        if technical_depth is not None:
            if not isinstance(technical_depth, (int, float)):
                errors.append("Technical depth must be a number")
            elif technical_depth < 1 or technical_depth > 5:
                errors.append("Technical depth must be between 1 and 5")
        
        # Interaction level validation
        valid_interactions = ["Minimal", "Moderate", "High"]
        interaction_level = params.get('interaction_level')
        if interaction_level and interaction_level not in valid_interactions:
            errors.append(f"Invalid interaction level. Must be one of: {', '.join(valid_interactions)}")
        
        # Language validation
        valid_languages = ["English", "Korean"]
        language = params.get('language')
        if language and language not in valid_languages:
            errors.append(f"Invalid language. Must be one of: {', '.join(valid_languages)}")
        
        is_valid = len(errors) == 0
        message = "Presentation parameters are valid" if is_valid else f"Found {len(errors)} validation errors"
        
        return ValidationResult(is_valid, message, errors)


class FileValidator:
    """Validates uploaded files and file operations."""
    
    @staticmethod
    def validate_pptx_file(file_path: str, max_size_mb: int = 50) -> ValidationResult:
        """Validate PowerPoint file.
        
        Args:
            file_path: Path to PowerPoint file
            max_size_mb: Maximum file size in MB
            
        Returns:
            ValidationResult object
        """
        errors = []
        
        try:
            path = Path(file_path)
            
            # Check if file exists
            if not path.exists():
                errors.append("File does not exist")
                return ValidationResult(False, "File validation failed", errors)
            
            # Check file extension
            if path.suffix.lower() != '.pptx':
                errors.append("File must be a PowerPoint (.pptx) file")
            
            # Check file size
            file_size_mb = path.stat().st_size / (1024 * 1024)
            if file_size_mb > max_size_mb:
                errors.append(f"File size ({file_size_mb:.1f} MB) exceeds maximum ({max_size_mb} MB)")
            
            # Check if file is readable
            if not path.is_file():
                errors.append("Path is not a valid file")
            
            # Try to open file to check if it's corrupted
            try:
                from pptx import Presentation
                pres = Presentation(file_path)
                if len(pres.slides) == 0:
                    errors.append("PowerPoint file contains no slides")
                elif len(pres.slides) > 100:
                    errors.append("PowerPoint file contains too many slides (maximum 100)")
            except Exception as e:
                errors.append(f"PowerPoint file appears to be corrupted: {str(e)}")
            
        except Exception as e:
            errors.append(f"File validation error: {str(e)}")
        
        is_valid = len(errors) == 0
        message = "PowerPoint file is valid" if is_valid else f"Found {len(errors)} validation errors"
        
        return ValidationResult(is_valid, message, errors)
    
    @staticmethod
    def validate_file_upload(uploaded_file: Any) -> ValidationResult:
        """Validate Streamlit uploaded file.
        
        Args:
            uploaded_file: Streamlit uploaded file object
            
        Returns:
            ValidationResult object
        """
        errors = []
        
        if uploaded_file is None:
            errors.append("No file uploaded")
            return ValidationResult(False, "No file provided", errors)
        
        # Check file name
        if not uploaded_file.name:
            errors.append("File has no name")
        elif not uploaded_file.name.endswith('.pptx'):
            errors.append("File must have .pptx extension")
        
        # Check file size
        file_size = len(uploaded_file.getvalue())
        file_size_mb = file_size / (1024 * 1024)
        
        if file_size == 0:
            errors.append("File is empty")
        elif file_size_mb > 50:
            errors.append(f"File size ({file_size_mb:.1f} MB) exceeds maximum (50 MB)")
        
        # Check file type
        if uploaded_file.type and not uploaded_file.type.startswith('application/vnd.openxmlformats'):
            errors.append(f"Invalid file type: {uploaded_file.type}")
        
        is_valid = len(errors) == 0
        message = "Uploaded file is valid" if is_valid else f"Found {len(errors)} validation errors"
        
        return ValidationResult(is_valid, message, errors)


class ContentValidator:
    """Validates generated content quality."""
    
    @staticmethod
    def validate_script_content(script: str, min_length: int = 100) -> ValidationResult:
        """Validate generated script content.
        
        Args:
            script: Generated script text
            min_length: Minimum script length in characters
            
        Returns:
            ValidationResult object
        """
        errors = []
        
        if not script or not script.strip():
            errors.append("Script is empty")
            return ValidationResult(False, "Script validation failed", errors)
        
        script = script.strip()
        
        # Check minimum length
        if len(script) < min_length:
            errors.append(f"Script is too short (minimum {min_length} characters)")
        
        # Check for basic structure
        if not re.search(r'#.*slide', script, re.IGNORECASE):
            errors.append("Script appears to be missing slide sections")
        
        # Check for time allocations
        if not re.search(r'\d+\s*minutes?', script, re.IGNORECASE):
            errors.append("Script appears to be missing time allocations")
        
        # Check for reasonable word count
        word_count = len(script.split())
        if word_count < 50:
            errors.append("Script has too few words")
        elif word_count > 10000:
            errors.append("Script is unusually long")
        
        is_valid = len(errors) == 0
        message = "Script content is valid" if is_valid else f"Found {len(errors)} validation errors"
        
        return ValidationResult(is_valid, message, errors)
    
    @staticmethod
    def validate_time_allocation(slides: List[Dict], total_duration: int) -> ValidationResult:
        """Validate time allocation across slides.
        
        Args:
            slides: List of slide data with time allocations
            total_duration: Total presentation duration in minutes
            
        Returns:
            ValidationResult object
        """
        errors = []
        
        if not slides:
            errors.append("No slides provided for time validation")
            return ValidationResult(False, "Time allocation validation failed", errors)
        
        # Calculate total allocated time
        total_allocated = 0
        for slide in slides:
            slide_time = slide.get('allocated_time', 0)
            if not isinstance(slide_time, (int, float)) or slide_time <= 0:
                errors.append(f"Invalid time allocation for slide {slide.get('number', '?')}")
            else:
                total_allocated += slide_time
        
        # Check if total time matches duration (with 10% tolerance)
        tolerance = total_duration * 0.1
        if abs(total_allocated - total_duration) > tolerance:
            errors.append(
                f"Total allocated time ({total_allocated:.1f} min) doesn't match "
                f"presentation duration ({total_duration} min)"
            )
        
        # Check for reasonable per-slide times
        avg_time = total_duration / len(slides)
        for slide in slides:
            slide_time = slide.get('allocated_time', 0)
            if slide_time > avg_time * 3:  # More than 3x average
                errors.append(f"Slide {slide.get('number', '?')} has unusually long time allocation")
            elif slide_time < avg_time * 0.2:  # Less than 20% of average
                errors.append(f"Slide {slide.get('number', '?')} has unusually short time allocation")
        
        is_valid = len(errors) == 0
        message = "Time allocation is valid" if is_valid else f"Found {len(errors)} validation errors"
        
        return ValidationResult(is_valid, message, errors)


def validate_all_inputs(persona_data: Dict, params: Dict, uploaded_file: Any) -> ValidationResult:
    """Validate all user inputs comprehensively.
    
    Args:
        persona_data: SA persona data
        params: Presentation parameters
        uploaded_file: Uploaded PowerPoint file
        
    Returns:
        Combined ValidationResult
    """
    all_errors = []
    
    # Validate persona data
    persona_result = PersonaValidator.validate_persona_data(persona_data)
    if not persona_result.is_valid:
        all_errors.extend([f"Persona: {error}" for error in persona_result.errors])
    
    # Validate presentation parameters
    params_result = PresentationValidator.validate_presentation_params(params)
    if not params_result.is_valid:
        all_errors.extend([f"Parameters: {error}" for error in params_result.errors])
    
    # Validate uploaded file
    file_result = FileValidator.validate_file_upload(uploaded_file)
    if not file_result.is_valid:
        all_errors.extend([f"File: {error}" for error in file_result.errors])
    
    is_valid = len(all_errors) == 0
    message = "All inputs are valid" if is_valid else f"Found {len(all_errors)} validation errors"
    
    return ValidationResult(is_valid, message, all_errors)
