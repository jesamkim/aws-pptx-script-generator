"""PowerPoint Processing Engine.

This module handles PowerPoint file processing, including slide extraction,
content analysis, and metadata collection.
"""

from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import tempfile
from dataclasses import dataclass
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from loguru import logger
import io


@dataclass
class SlideMetadata:
    """Metadata for a single slide.
    
    Attributes:
        slide_number: Slide number (1-based)
        title: Slide title text
        layout_name: Slide layout name
        shape_count: Number of shapes on slide
        has_notes: Whether slide has speaker notes
        has_images: Whether slide contains images
        has_charts: Whether slide contains charts
        has_tables: Whether slide contains tables
    """
    slide_number: int
    title: str
    layout_name: str
    shape_count: int
    has_notes: bool
    has_images: bool
    has_charts: bool
    has_tables: bool


@dataclass
class SlideContent:
    """Content extracted from a single slide.
    
    Attributes:
        metadata: Slide metadata
        text_content: All text content from slide
        speaker_notes: Speaker notes text
        images: List of image information
        charts: List of chart information
        tables: List of table data
        hyperlinks: List of hyperlinks found
    """
    metadata: SlideMetadata
    text_content: List[str]
    speaker_notes: str
    images: List[Dict[str, Any]]
    charts: List[Dict[str, Any]]
    tables: List[Dict[str, Any]]
    hyperlinks: List[Dict[str, str]]


@dataclass
class PresentationData:
    """Complete presentation data.
    
    Attributes:
        title: Presentation title
        slide_count: Total number of slides
        slides: List of slide content
        metadata: Presentation-level metadata
        file_path: Path to original file
        file_size: File size in bytes
    """
    title: str
    slide_count: int
    slides: List[SlideContent]
    metadata: Dict[str, Any]
    file_path: str
    file_size: int


class PowerPointProcessor:
    """PowerPoint file processor with comprehensive content extraction.
    
    This class handles loading, parsing, and extracting content from PowerPoint
    presentations with robust error handling and validation.
    """
    
    def __init__(self):
        """Initialize PowerPoint processor."""
        self.presentation: Optional[Presentation] = None
        self.file_path: Optional[str] = None
        logger.info("Initialized PowerPoint processor")
    
    def load_presentation(self, file_path: str) -> bool:
        """Load PowerPoint presentation from file path.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            True if loaded successfully, False otherwise
            
        Raises:
            Exception: If file loading fails
        """
        try:
            self.file_path = file_path
            self.presentation = Presentation(file_path)
            logger.info(f"Successfully loaded presentation: {file_path}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to load presentation {file_path}: {str(e)}")
            raise Exception(f"Could not load PowerPoint file: {str(e)}")
    
    def load_from_bytes(self, file_bytes: bytes, filename: str) -> bool:
        """Load PowerPoint presentation from bytes.
        
        Args:
            file_bytes: PowerPoint file as bytes
            filename: Original filename for reference
            
        Returns:
            True if loaded successfully, False otherwise
            
        Raises:
            Exception: If file loading fails
        """
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                tmp_file.write(file_bytes)
                tmp_file.flush()
                
                self.file_path = tmp_file.name
                self.presentation = Presentation(io.BytesIO(file_bytes))
                
            logger.info(f"Successfully loaded presentation from bytes: {filename}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to load presentation from bytes {filename}: {str(e)}")
            raise Exception(f"Could not load PowerPoint file: {str(e)}")
    
    def validate_file_integrity(self) -> Dict[str, Any]:
        """Validate PowerPoint file integrity and structure.
        
        Returns:
            Validation results dictionary
            
        Raises:
            Exception: If validation fails
        """
        if not self.presentation:
            raise Exception("No presentation loaded")
        
        try:
            validation_result = {
                'valid': True,
                'issues': [],
                'warnings': [],
                'slide_count': len(self.presentation.slides),
                'has_master_slides': len(self.presentation.slide_masters) > 0,
                'has_layouts': len(self.presentation.slide_layouts) > 0,
            }
            
            # Check for empty presentation
            if validation_result['slide_count'] == 0:
                validation_result['valid'] = False
                validation_result['issues'].append("Presentation contains no slides")
            
            # Check for corrupted slides
            corrupted_slides = []
            for i, slide in enumerate(self.presentation.slides):
                try:
                    # Try to access slide properties
                    _ = slide.shapes
                    _ = slide.slide_layout
                except Exception as e:
                    corrupted_slides.append(i + 1)
                    validation_result['warnings'].append(f"Slide {i + 1} may be corrupted: {str(e)}")
            
            if corrupted_slides:
                validation_result['corrupted_slides'] = corrupted_slides
            
            # Check file size limits
            if self.file_path and Path(self.file_path).stat().st_size > 50 * 1024 * 1024:  # 50MB
                validation_result['warnings'].append("File size exceeds recommended 50MB limit")
            
            logger.info(f"File validation completed: {validation_result['valid']}")
            return validation_result
            
        except Exception as e:
            logger.error(f"File validation failed: {str(e)}")
            raise Exception(f"File validation error: {str(e)}")
    
    def extract_slide_metadata(self, slide: Slide, slide_number: int) -> SlideMetadata:
        """Extract metadata from a single slide.
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number (1-based)
            
        Returns:
            SlideMetadata object
        """
        try:
            # Extract title
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    title = shape.text.strip()
                    break
            
            # Count different shape types
            has_images = any(shape.shape_type == 13 for shape in slide.shapes)  # MSO_SHAPE_TYPE.PICTURE
            has_charts = any(shape.shape_type == 3 for shape in slide.shapes)   # MSO_SHAPE_TYPE.CHART
            has_tables = any(shape.shape_type == 19 for shape in slide.shapes)  # MSO_SHAPE_TYPE.TABLE
            
            # Check for speaker notes
            has_notes = bool(slide.notes_slide.notes_text_frame.text.strip()) if slide.notes_slide else False
            
            metadata = SlideMetadata(
                slide_number=slide_number,
                title=title[:100],  # Limit title length
                layout_name=slide.slide_layout.name if slide.slide_layout else "Unknown",
                shape_count=len(slide.shapes),
                has_notes=has_notes,
                has_images=has_images,
                has_charts=has_charts,
                has_tables=has_tables
            )
            
            logger.debug(f"Extracted metadata for slide {slide_number}")
            return metadata
            
        except Exception as e:
            logger.warning(f"Failed to extract metadata for slide {slide_number}: {str(e)}")
            return SlideMetadata(
                slide_number=slide_number,
                title="Error extracting title",
                layout_name="Unknown",
                shape_count=0,
                has_notes=False,
                has_images=False,
                has_charts=False,
                has_tables=False
            )
    
    def extract_text_content(self, slide: Slide) -> List[str]:
        """Extract all text content from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            List of text strings found on slide
        """
        text_content = []
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_content.append(shape.text.strip())
                elif hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_content.append(paragraph.text.strip())
            
            logger.debug(f"Extracted {len(text_content)} text elements from slide")
            return text_content
            
        except Exception as e:
            logger.warning(f"Failed to extract text content: {str(e)}")
            return []
    
    def extract_speaker_notes(self, slide: Slide) -> str:
        """Extract speaker notes from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Speaker notes text
        """
        try:
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                logger.debug(f"Extracted speaker notes: {len(notes_text)} characters")
                return notes_text
            return ""
            
        except Exception as e:
            logger.warning(f"Failed to extract speaker notes: {str(e)}")
            return ""
    
    def extract_hyperlinks(self, slide: Slide) -> List[Dict[str, str]]:
        """Extract hyperlinks from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            List of hyperlink dictionaries
        """
        hyperlinks = []
        
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.hyperlink and run.hyperlink.address:
                                hyperlinks.append({
                                    'text': run.text,
                                    'url': run.hyperlink.address,
                                    'type': 'external' if run.hyperlink.address.startswith('http') else 'internal'
                                })
            
            logger.debug(f"Extracted {len(hyperlinks)} hyperlinks from slide")
            return hyperlinks
            
        except Exception as e:
            logger.warning(f"Failed to extract hyperlinks: {str(e)}")
            return []
    
    def extract_slide_content(self, slide: Slide, slide_number: int) -> SlideContent:
        """Extract complete content from a single slide.
        
        Args:
            slide: PowerPoint slide object
            slide_number: Slide number (1-based)
            
        Returns:
            SlideContent object with all extracted data
        """
        try:
            metadata = self.extract_slide_metadata(slide, slide_number)
            text_content = self.extract_text_content(slide)
            speaker_notes = self.extract_speaker_notes(slide)
            hyperlinks = self.extract_hyperlinks(slide)
            
            # Placeholder for image, chart, and table extraction
            # These would be implemented with more detailed shape analysis
            images = []
            charts = []
            tables = []
            
            slide_content = SlideContent(
                metadata=metadata,
                text_content=text_content,
                speaker_notes=speaker_notes,
                images=images,
                charts=charts,
                tables=tables,
                hyperlinks=hyperlinks
            )
            
            logger.info(f"Successfully extracted content from slide {slide_number}")
            return slide_content
            
        except Exception as e:
            logger.error(f"Failed to extract content from slide {slide_number}: {str(e)}")
            raise Exception(f"Slide content extraction failed: {str(e)}")
    
    def process_presentation(self) -> PresentationData:
        """Process entire presentation and extract all content.
        
        Returns:
            PresentationData object with complete presentation information
            
        Raises:
            Exception: If processing fails
        """
        if not self.presentation:
            raise Exception("No presentation loaded")
        
        try:
            slides_content = []
            
            # Process each slide
            for i, slide in enumerate(self.presentation.slides, 1):
                slide_content = self.extract_slide_content(slide, i)
                slides_content.append(slide_content)
            
            # Extract presentation-level metadata
            presentation_title = ""
            if slides_content and slides_content[0].text_content:
                presentation_title = slides_content[0].text_content[0]
            
            # Get file size
            file_size = 0
            if self.file_path and Path(self.file_path).exists():
                file_size = Path(self.file_path).stat().st_size
            
            presentation_data = PresentationData(
                title=presentation_title[:200],  # Limit title length
                slide_count=len(slides_content),
                slides=slides_content,
                metadata={
                    'slide_masters': len(self.presentation.slide_masters),
                    'slide_layouts': len(self.presentation.slide_layouts),
                    'core_properties': {
                        'title': self.presentation.core_properties.title or "",
                        'author': self.presentation.core_properties.author or "",
                        'subject': self.presentation.core_properties.subject or "",
                        'created': str(self.presentation.core_properties.created) if self.presentation.core_properties.created else "",
                        'modified': str(self.presentation.core_properties.modified) if self.presentation.core_properties.modified else "",
                    }
                },
                file_path=self.file_path or "",
                file_size=file_size
            )
            
            logger.info(f"Successfully processed presentation: {len(slides_content)} slides")
            return presentation_data
            
        except Exception as e:
            logger.error(f"Failed to process presentation: {str(e)}")
            raise Exception(f"Presentation processing failed: {str(e)}")
    
    def update_speaker_notes(self, slide_number: int, notes_text: str) -> bool:
        """Update speaker notes for a specific slide.
        
        Args:
            slide_number: Slide number (1-based)
            notes_text: New speaker notes text
            
        Returns:
            True if update successful, False otherwise
        """
        if not self.presentation:
            logger.error("No presentation loaded")
            return False
        
        try:
            if slide_number < 1 or slide_number > len(self.presentation.slides):
                logger.error(f"Invalid slide number: {slide_number}")
                return False
            
            slide = self.presentation.slides[slide_number - 1]
            
            # Ensure notes slide exists
            if not slide.notes_slide:
                logger.warning(f"No notes slide for slide {slide_number}")
                return False
            
            # Update notes text
            slide.notes_slide.notes_text_frame.clear()
            slide.notes_slide.notes_text_frame.text = notes_text
            
            logger.info(f"Updated speaker notes for slide {slide_number}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to update speaker notes for slide {slide_number}: {str(e)}")
            return False
    
    def save_presentation(self, output_path: str) -> bool:
        """Save presentation to file.
        
        Args:
            output_path: Path for output file
            
        Returns:
            True if save successful, False otherwise
        """
        if not self.presentation:
            logger.error("No presentation loaded")
            return False
        
        try:
            self.presentation.save(output_path)
            logger.info(f"Saved presentation to: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to save presentation to {output_path}: {str(e)}")
            return False
