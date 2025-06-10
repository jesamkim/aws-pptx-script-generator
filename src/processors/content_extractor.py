"""Content Extraction Module.

This module handles structured content extraction from PowerPoint slides,
including text, images, charts, tables, and AWS service references.
"""

from typing import List, Dict, Any, Optional, Set
from dataclasses import dataclass
import re
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.chart import Chart
from pptx.table import Table
from loguru import logger


@dataclass
class TextElement:
    """Structured text element from a slide.
    
    Attributes:
        text: Text content
        level: Hierarchy level (0 = title, 1+ = bullet points)
        font_size: Font size in points
        is_bold: Whether text is bold
        is_italic: Whether text is italic
        color: Text color (RGB hex)
        position: Position on slide (x, y coordinates)
    """
    text: str
    level: int
    font_size: Optional[float] = None
    is_bold: bool = False
    is_italic: bool = False
    color: str = "#000000"
    position: Optional[Dict[str, float]] = None


@dataclass
class ImageElement:
    """Image element from a slide.
    
    Attributes:
        description: Image description or alt text
        file_type: Image file type (e.g., PNG, JPEG)
        size: Image dimensions (width, height)
        position: Position on slide
        data: Image data (if extracted)
    """
    description: str
    file_type: str
    size: Dict[str, int]
    position: Dict[str, float]
    data: Optional[bytes] = None


@dataclass
class ChartElement:
    """Chart element from a slide.
    
    Attributes:
        chart_type: Type of chart (e.g., bar, line, pie)
        title: Chart title
        data: Chart data as dictionary
        position: Position on slide
    """
    chart_type: str
    title: str
    data: Dict[str, Any]
    position: Dict[str, float]


@dataclass
class TableElement:
    """Table element from a slide.
    
    Attributes:
        rows: Number of rows
        columns: Number of columns
        headers: Table headers
        data: Table data as list of rows
        position: Position on slide
    """
    rows: int
    columns: int
    headers: List[str]
    data: List[List[str]]
    position: Dict[str, float]


@dataclass
class AWSService:
    """AWS service reference found in content.
    
    Attributes:
        name: Service name (e.g., "Amazon S3", "AWS Lambda")
        short_name: Short service name (e.g., "S3", "Lambda")
        mentions: Number of mentions in content
        context: List of text snippets containing service mentions
    """
    name: str
    short_name: str
    mentions: int
    context: List[str]


class ContentExtractor:
    """Extracts and structures content from PowerPoint slides.
    
    This class provides comprehensive content extraction capabilities,
    including text hierarchy, images, charts, tables, and AWS service
    identification.
    """
    
    def __init__(self):
        """Initialize content extractor with AWS service patterns."""
        # Common AWS service name patterns
        self.aws_patterns = {
            r"Amazon\s+S3": "S3",
            r"Amazon\s+EC2": "EC2",
            r"AWS\s+Lambda": "Lambda",
            r"Amazon\s+RDS": "RDS",
            r"Amazon\s+DynamoDB": "DynamoDB",
            r"AWS\s+CloudFormation": "CloudFormation",
            r"Amazon\s+VPC": "VPC",
            r"Amazon\s+CloudWatch": "CloudWatch",
            r"AWS\s+IAM": "IAM",
            r"Amazon\s+Route\s*53": "Route 53",
            r"Amazon\s+SNS": "SNS",
            r"Amazon\s+SQS": "SQS",
            r"AWS\s+Step\s+Functions": "Step Functions",
            r"Amazon\s+ECS": "ECS",
            r"Amazon\s+EKS": "EKS",
            r"Amazon\s+Aurora": "Aurora",
            r"AWS\s+Glue": "Glue",
            r"Amazon\s+EMR": "EMR",
            r"Amazon\s+Redshift": "Redshift",
            r"AWS\s+Organizations": "Organizations",
            # Add more patterns as needed
        }
        
        logger.info("Initialized content extractor with AWS service patterns")
    
    def extract_text_hierarchy(self, shape: BaseShape) -> List[TextElement]:
        """Extract text with hierarchy information from a shape.
        
        Args:
            shape: PowerPoint shape object
            
        Returns:
            List of TextElement objects
        """
        text_elements = []
        
        try:
            if not hasattr(shape, 'text_frame'):
                return text_elements
            
            for paragraph in shape.text_frame.paragraphs:
                # Get paragraph level (0 = title, 1+ = bullet points)
                level = paragraph.level
                
                # Get font properties from first run (if any)
                font_props = {}
                if paragraph.runs:
                    font = paragraph.runs[0].font
                    font_props = {
                        'size': font.size.pt if font.size else None,
                        'bold': font.bold if hasattr(font, 'bold') else False,
                        'italic': font.italic if hasattr(font, 'italic') else False,
                        'color': f"#{font.color.rgb:06x}" if font.color and font.color.rgb else "#000000"
                    }
                
                # Get shape position
                position = {
                    'x': shape.left,
                    'y': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }
                
                text_element = TextElement(
                    text=paragraph.text.strip(),
                    level=level,
                    font_size=font_props.get('size'),
                    is_bold=font_props.get('bold', False),
                    is_italic=font_props.get('italic', False),
                    color=font_props.get('color', "#000000"),
                    position=position
                )
                
                if text_element.text:  # Only add non-empty text
                    text_elements.append(text_element)
            
            logger.debug(f"Extracted {len(text_elements)} text elements from shape")
            return text_elements
            
        except Exception as e:
            logger.warning(f"Failed to extract text hierarchy: {str(e)}")
            return text_elements
    
    def extract_image_info(self, picture: Picture) -> Optional[ImageElement]:
        """Extract information about an image shape.
        
        Args:
            picture: PowerPoint picture shape
            
        Returns:
            ImageElement object or None if extraction fails
        """
        try:
            # Get image description/alt text
            description = picture.alt_text if hasattr(picture, 'alt_text') else ""
            
            # Get image file type
            file_type = "Unknown"
            if hasattr(picture, 'image') and hasattr(picture.image, 'ext'):
                file_type = picture.image.ext.upper()
            
            # Get image size and position
            size = {
                'width': picture.width,
                'height': picture.height
            }
            
            position = {
                'x': picture.left,
                'y': picture.top,
                'width': picture.width,
                'height': picture.height
            }
            
            # Create image element
            image_element = ImageElement(
                description=description,
                file_type=file_type,
                size=size,
                position=position
            )
            
            logger.debug(f"Extracted image info: {file_type}, {size}")
            return image_element
            
        except Exception as e:
            logger.warning(f"Failed to extract image info: {str(e)}")
            return None
    
    def extract_chart_data(self, chart: Chart) -> Optional[ChartElement]:
        """Extract data from a chart shape.
        
        Args:
            chart: PowerPoint chart object
            
        Returns:
            ChartElement object or None if extraction fails
        """
        try:
            # Get chart type
            chart_type = chart.chart_type.name if hasattr(chart, 'chart_type') else "Unknown"
            
            # Get chart title
            title = ""
            if hasattr(chart, 'chart_title') and chart.chart_title:
                title = chart.chart_title.text_frame.text
            
            # Get chart position
            position = {
                'x': chart.left,
                'y': chart.top,
                'width': chart.width,
                'height': chart.height
            }
            
            # Extract chart data (simplified)
            data = {
                'series': [],
                'categories': [],
                'values': []
            }
            
            # Create chart element
            chart_element = ChartElement(
                chart_type=chart_type,
                title=title,
                data=data,
                position=position
            )
            
            logger.debug(f"Extracted chart data: {chart_type}, {title}")
            return chart_element
            
        except Exception as e:
            logger.warning(f"Failed to extract chart data: {str(e)}")
            return None
    
    def extract_table_data(self, table: Table) -> Optional[TableElement]:
        """Extract data from a table shape.
        
        Args:
            table: PowerPoint table object
            
        Returns:
            TableElement object or None if extraction fails
        """
        try:
            rows = len(table.rows)
            columns = len(table.columns)
            
            # Extract headers (first row)
            headers = []
            for cell in table.rows[0].cells:
                headers.append(cell.text.strip())
            
            # Extract data (remaining rows)
            data = []
            for row in table.rows[1:]:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                data.append(row_data)
            
            # Get table position
            position = {
                'x': table.left,
                'y': table.top,
                'width': table.width,
                'height': table.height
            }
            
            # Create table element
            table_element = TableElement(
                rows=rows,
                columns=columns,
                headers=headers,
                data=data,
                position=position
            )
            
            logger.debug(f"Extracted table data: {rows}x{columns}")
            return table_element
            
        except Exception as e:
            logger.warning(f"Failed to extract table data: {str(e)}")
            return None
    
    def identify_aws_services(self, text: str) -> List[AWSService]:
        """Identify AWS service references in text.
        
        Args:
            text: Text content to analyze
            
        Returns:
            List of AWSService objects found
        """
        services = {}
        
        try:
            # Search for each service pattern
            for pattern, short_name in self.aws_patterns.items():
                matches = re.finditer(pattern, text, re.IGNORECASE)
                
                for match in matches:
                    service_name = match.group(0)
                    
                    # Get context (text around the mention)
                    start = max(0, match.start() - 50)
                    end = min(len(text), match.end() + 50)
                    context = text[start:end].strip()
                    
                    # Update or create service entry
                    if service_name in services:
                        services[service_name].mentions += 1
                        services[service_name].context.append(context)
                    else:
                        services[service_name] = AWSService(
                            name=service_name,
                            short_name=short_name,
                            mentions=1,
                            context=[context]
                        )
            
            logger.debug(f"Identified {len(services)} AWS services in text")
            return list(services.values())
            
        except Exception as e:
            logger.warning(f"Failed to identify AWS services: {str(e)}")
            return []
    
    def extract_slide_content(self, slide: Any) -> Dict[str, Any]:
        """Extract all content from a slide.
        
        Args:
            slide: PowerPoint slide object
            
        Returns:
            Dictionary containing all extracted content
        """
        content = {
            'text_elements': [],
            'images': [],
            'charts': [],
            'tables': [],
            'aws_services': [],
            'all_text': ""
        }
        
        try:
            all_text = []
            
            for shape in slide.shapes:
                # Extract text content
                if hasattr(shape, 'text_frame'):
                    text_elements = self.extract_text_hierarchy(shape)
                    content['text_elements'].extend(text_elements)
                    all_text.extend(element.text for element in text_elements)
                
                # Extract images
                if isinstance(shape, Picture):
                    image_info = self.extract_image_info(shape)
                    if image_info:
                        content['images'].append(image_info)
                
                # Extract charts
                if isinstance(shape, GraphicFrame) and shape.chart:
                    chart_data = self.extract_chart_data(shape.chart)
                    if chart_data:
                        content['charts'].append(chart_data)
                
                # Extract tables
                if isinstance(shape, GraphicFrame) and shape.table:
                    table_data = self.extract_table_data(shape.table)
                    if table_data:
                        content['tables'].append(table_data)
            
            # Combine all text for AWS service identification
            content['all_text'] = ' '.join(all_text)
            content['aws_services'] = self.identify_aws_services(content['all_text'])
            
            logger.info(f"Extracted complete slide content: {len(content['text_elements'])} text elements, "
                       f"{len(content['images'])} images, {len(content['charts'])} charts, "
                       f"{len(content['tables'])} tables, {len(content['aws_services'])} AWS services")
            return content
            
        except Exception as e:
            logger.error(f"Failed to extract slide content: {str(e)}")
            return content
    
    def extract_presentation_content(self, slides: List[Any]) -> List[Dict[str, Any]]:
        """Extract content from all slides in a presentation.
        
        Args:
            slides: List of PowerPoint slide objects
            
        Returns:
            List of content dictionaries for each slide
        """
        content_list = []
        
        try:
            for i, slide in enumerate(slides):
                logger.info(f"Processing slide {i + 1}/{len(slides)}")
                content = self.extract_slide_content(slide)
                content_list.append(content)
            
            logger.info(f"Extracted content from {len(content_list)} slides")
            return content_list
            
        except Exception as e:
            logger.error(f"Failed to extract presentation content: {str(e)}")
            return content_list
    
    def get_aws_service_summary(self, content_list: List[Dict[str, Any]]) -> Dict[str, AWSService]:
        """Generate summary of AWS services across all slides.
        
        Args:
            content_list: List of slide content dictionaries
            
        Returns:
            Dictionary of AWSService objects with total mentions
        """
        service_summary = {}
        
        try:
            for content in content_list:
                for service in content['aws_services']:
                    if service.name in service_summary:
                        service_summary[service.name].mentions += service.mentions
                        service_summary[service.name].context.extend(service.context)
                    else:
                        service_summary[service.name] = service
            
            logger.info(f"Generated AWS service summary: {len(service_summary)} unique services")
            return service_summary
            
        except Exception as e:
            logger.error(f"Failed to generate AWS service summary: {str(e)}")
            return service_summary
